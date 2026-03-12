"""生成徐记海鲜POC汇报PPT (.pptx) — 15页"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = os.path.join(os.path.dirname(__file__), "智链经营助手_徐记海鲜_POC汇报.pptx")

# ── 品牌色 ──
BRAND = RGBColor(0xFF, 0x6B, 0x2C)  # #FF6B2C
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xFE, 0xF5, 0xF0)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)


def add_brand_bar(slide, prs):
    """左侧品牌色竖条"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(0.15), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND
    shape.line.fill.background()


def add_footer(slide, page_num, total=15):
    """底部页码"""
    txBox = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(1.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def set_text(shape, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=14, bold=False, color=DARK, bullet=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if bullet:
        p.level = 0
    return p


def make_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    # 品牌色表头
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.size = Pt(10)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    return table


# ═══════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════

def slide_01_cover(prs):
    """封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 大色块背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()
    # 品牌色横条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.0), Inches(8), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()
    # 标题
    title = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    set_text(title, "智链经营助手 × 徐记海鲜", 36, True, WHITE, PP_ALIGN.LEFT)
    add_paragraph(title.text_frame, "POC验证报告", 28, False, BRAND)
    # 副标题
    sub = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(8), Inches(0.8))
    set_text(sub, "AI驱动的连锁餐饮经营决策系统", 16, False, RGBColor(0xAA, 0xAA, 0xAA))
    add_paragraph(sub.text_frame, "屯象科技  |  2026年3月", 12, False, GRAY)


def slide_02_profile(prs):
    """徐记海鲜企业画像"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "徐记海鲜 — 企业画像", 28, True, DARK)

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    items = [
        ("品牌", "徐记海鲜（2018年创立，长沙总部）"),
        ("规模", "5家门店（长沙3 + 株洲1 + 湘潭1）"),
        ("年营收", "~¥8,000万"),
        ("旗舰店", "五一广场店（800㎡，45人，日均¥6-8万）"),
    ]
    for label, val in items:
        p = tf.add_paragraph()
        p.text = f"  {label}：{val}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)

    # 痛点框
    pain = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(2.2))
    pain.fill.solid()
    pain.fill.fore_color.rgb = LIGHT_BG
    pain.line.color.rgb = BRAND
    ptf = pain.text_frame
    ptf.word_wrap = True
    p0 = ptf.paragraphs[0]
    p0.text = "核心痛点"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    pains = [
        "食材成本率偏高：行业28-33%，徐记实际~35%（差距7个点）",
        "人力排班低效：凭经验排班，高峰缺人/低谷冗余",
        "跨店管理不透明：5店数据割裂，总部无法实时对标",
    ]
    for pt in pains:
        pp = ptf.add_paragraph()
        pp.text = f"  ⚠ {pt}"
        pp.font.size = Pt(12)
        pp.font.color.rgb = DARK

    add_footer(slide, 2)


def slide_03_value(prs):
    """核心价值主张"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "核心价值主张", 28, True, DARK)

    # 大数字
    big = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    set_text(big, "帮每家连锁店每年多赚 ¥30万+", 32, True, BRAND, PP_ALIGN.CENTER)
    add_paragraph(big.text_frame, "成本率降低2个点 → 年均节省¥30万+", 16, False, GRAY, PP_ALIGN.CENTER)

    # 三个"一"
    cards = [
        ("一个助手", "企业微信", "4时间点智能推送\n一键确认决策"),
        ("一个看板", "BFF聚合", "30秒刷新\n4角色定制视图"),
        ("一套大脑", "11 AI Agent", "三级降级\n自动学习进化"),
    ]
    for i, (t, sub, desc) in enumerate(cards):
        left = Inches(0.5 + i * 3.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.5), Inches(2.8), Inches(3.0))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG if i != 1 else BRAND
        box.line.color.rgb = BRAND
        btf = box.text_frame
        btf.word_wrap = True
        p0 = btf.paragraphs[0]
        p0.text = t
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = DARK if i != 1 else WHITE
        p0.alignment = PP_ALIGN.CENTER
        p1 = btf.add_paragraph()
        p1.text = sub
        p1.font.size = Pt(14)
        p1.font.color.rgb = BRAND if i != 1 else WHITE
        p1.alignment = PP_ALIGN.CENTER
        p2 = btf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY if i != 1 else RGBColor(0xFF, 0xDD, 0xCC)
        p2.alignment = PP_ALIGN.CENTER

    # 北极星
    ns = slide.shapes.add_textbox(Inches(3), Inches(6.8), Inches(4), Inches(0.4))
    set_text(ns, "北极星指标：续费率 ≥ 95%", 12, True, GRAY, PP_ALIGN.CENTER)

    add_footer(slide, 3)


def slide_04_architecture(prs):
    """四层架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "四层 Ontology 架构", 28, True, DARK)

    layers = [
        ("感知层 Perception", "企微 / 天财商龙POS / 飞书 数据采集", RGBColor(0x34, 0x98, 0xDB)),
        ("本体层 Ontology", "62个数据模型，97个 Alembic 迁移", RGBColor(0x27, 0xAE, 0x60)),
        ("推理层 Reasoning", "11个 AI Agent，三级降级策略", BRAND),
        ("行动层 Action", "企微推送 + 一键确认 + FSM 状态机", RGBColor(0x8E, 0x44, 0xAD)),
    ]
    for i, (name, desc, color) in enumerate(layers):
        top = Inches(1.3 + i * 1.4)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), top, Inches(8), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        btf = box.text_frame
        btf.word_wrap = True
        p0 = btf.paragraphs[0]
        p0.text = name
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = WHITE
        p0.alignment = PP_ALIGN.CENTER
        p1 = btf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xEE)
        p1.alignment = PP_ALIGN.CENTER

        # 箭头（除最后一层）
        if i < 3:
            arrow = slide.shapes.add_textbox(Inches(4.5), Inches(1.3 + i * 1.4 + 1.1), Inches(1), Inches(0.3))
            set_text(arrow, "▼", 16, True, GRAY, PP_ALIGN.CENTER)

    # 关键数字
    kpi = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
    set_text(kpi, "11个Agent × 4个决策窗口 × 30秒BFF刷新", 12, True, BRAND, PP_ALIGN.CENTER)

    add_footer(slide, 4)


def slide_05_mvp(prs):
    """10个MVP功能"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "10个MVP功能 — 全部就绪", 28, True, DARK)

    headers = ["#", "功能", "状态", "¥价值"]
    rows_data = [
        ["1", "POS数据自动采集", "✅", "日省1小时"],
        ["2", "每日利润快报", "✅", "实时经营状况"],
        ["3", "损耗Top5排名", "✅", "年省¥5-10万"],
        ["4", "决策型企微推送", "✅", "4时间点精准推送"],
        ["5", "一键审批采购建议", "✅", "采购效率+30%"],
        ["6", "BOM配方管理", "✅", "标准化成本核算"],
        ["7", "成本率趋势图", "✅", "预警前置"],
        ["8", "异常告警推送", "✅", "≤30min触达"],
        ["9", "月度经营报告", "✅", "一键生成"],
        ["10", "离线基础查询", "✅", "断网不停服"],
    ]
    table = make_table(slide, 11, 4, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(rows_data, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
            if c == 2:  # 状态列绿色
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = GREEN

    # 列宽
    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(1.0)
    table.columns[3].width = Inches(3.5)

    add_footer(slide, 5)


def slide_06_food_cost(prs):
    """食材成本分析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "食材成本分析 — 核心卖点", 28, True, DARK)

    # 对比数字
    boxes = [
        ("理论成本率", "28%", "(BOM标准)", GREEN),
        ("实际成本率", "35%", "(采购+损耗)", RED),
        ("差异", "7个点", "年损失¥56万", BRAND),
    ]
    for i, (label, num, sub, color) in enumerate(boxes):
        left = Inches(0.5 + i * 3.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(2.8), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = color
        btf = box.text_frame
        btf.word_wrap = True
        p0 = btf.paragraphs[0]
        p0.text = label
        p0.font.size = Pt(12)
        p0.font.color.rgb = GRAY
        p0.alignment = PP_ALIGN.CENTER
        p1 = btf.add_paragraph()
        p1.text = num
        p1.font.size = Pt(32)
        p1.font.bold = True
        p1.font.color.rgb = color
        p1.alignment = PP_ALIGN.CENTER
        p2 = btf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    # BOM闭环
    flow = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.6))
    set_text(flow, "BOM管理闭环：配方录入 → 成本计算 → 差异分析 → 决策建议 → 一键采购", 14, True, DARK, PP_ALIGN.CENTER)

    # Top5
    content = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(2.5))
    tf = content.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "Top5成本偏差食材"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    items = [
        "海鲜类：虾、蟹、鱼 — 供应商价格波动大，占总成本60%+",
        "酱料类：定制酱料 — 配比不标准，各店差异达20%",
        "蔬菜类：时令蔬菜 — 损耗率高达15-25%",
    ]
    for item in items:
        pp = tf.add_paragraph()
        pp.text = f"  • {item}"
        pp.font.size = Pt(12)
        pp.font.color.rgb = DARK

    add_footer(slide, 6)


def slide_07_health(prs):
    """门店健康指数"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "门店健康指数 — 5维评估", 28, True, DARK)

    headers = ["维度", "权重", "计算方式", "五一店模拟分"]
    data = [
        ["营收完成率", "30%", "实际/目标", "85"],
        ["翻台率", "20%", "实际/2.0标准", "72"],
        ["成本率", "25%", "FoodCost方差", "60 ⚠"],
        ["投诉率", "15%", "质检不合格/总数", "88"],
        ["员工效率", "10%", "人均产值/¥500", "78"],
    ]
    table = make_table(slide, 6, 4, Inches(0.5), Inches(1.3), Inches(9), Inches(3.2))
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                if r == 3 and c == 3:  # 成本率警告
                    p.font.color.rgb = RED

    # 综合分数
    score = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.0), Inches(5), Inches(1.8))
    score.fill.solid()
    score.fill.fore_color.rgb = LIGHT_BG
    score.line.color.rgb = BRAND
    stf = score.text_frame
    stf.word_wrap = True
    p0 = stf.paragraphs[0]
    p0.text = '综合分数：76.1 / 100 → "良好"'
    p0.font.size = Pt(20)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    p0.alignment = PP_ALIGN.CENTER
    p1 = stf.add_paragraph()
    p1.text = "关键洞察：成本率是最大短板（仅60分），优化空间最大"
    p1.font.size = Pt(13)
    p1.font.color.rgb = DARK
    p1.alignment = PP_ALIGN.CENTER

    add_footer(slide, 7)


def slide_08_decision(prs):
    """智能决策引擎"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "智能决策引擎", 28, True, DARK)

    # 公式
    formula = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(1.0))
    formula.fill.solid()
    formula.fill.fore_color.rgb = DARK
    formula.line.fill.background()
    ftf = formula.text_frame
    ftf.word_wrap = True
    p0 = ftf.paragraphs[0]
    p0.text = "得分 = 0.40×财务影响 + 0.30×紧急程度 + 0.20×置信度 + 0.10×执行难度"
    p0.font.size = Pt(14)
    p0.font.color.rgb = BRAND
    p0.font.bold = True
    p0.alignment = PP_ALIGN.CENTER

    # 4个决策窗口
    windows = [
        ("08:00", "晨推", "今日采购建议"),
        ("12:00", "午推", "午市实时调整"),
        ("17:30", "战前", "晚市备货/排班"),
        ("20:30", "晚推", "当日复盘/次日预警"),
    ]
    for i, (time, name, desc) in enumerate(windows):
        left = Inches(0.3 + i * 2.45)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.6), Inches(2.2), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = BRAND if i == 2 else LIGHT_BG
        box.line.color.rgb = BRAND
        btf = box.text_frame
        btf.word_wrap = True
        p0 = btf.paragraphs[0]
        p0.text = time
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = WHITE if i == 2 else BRAND
        p0.alignment = PP_ALIGN.CENTER
        p1 = btf.add_paragraph()
        p1.text = name
        p1.font.size = Pt(14)
        p1.font.color.rgb = WHITE if i == 2 else DARK
        p1.alignment = PP_ALIGN.CENTER
        p2 = btf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0xFF, 0xDD, 0xCC) if i == 2 else GRAY
        p2.alignment = PP_ALIGN.CENTER

    # 决策样例
    sample = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.6), Inches(9), Inches(2.2))
    sample.fill.solid()
    sample.fill.fore_color.rgb = LIGHT_BG
    sample.line.color.rgb = BRAND
    stf = sample.text_frame
    stf.word_wrap = True
    p0 = stf.paragraphs[0]
    p0.text = "决策推送样例"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    p1 = stf.add_paragraph()
    p1.text = '"建议减少明日虾采购量20%（库存充足），预期节省¥340，置信度82%"'
    p1.font.size = Pt(13)
    p1.font.color.rgb = DARK
    p2 = stf.add_paragraph()
    p2.text = "  [一键确认]    [查看详情]    [忽略]"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = BRAND

    add_footer(slide, 8)


def slide_09_wechat(prs):
    """企业微信集成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "企业微信集成", 28, True, DARK)

    # 左：消息模板
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(3.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "4种消息模板"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    templates = [
        "折扣审批：金额+操作人+订单号",
        "异常告警：类型+严重度+描述",
        "班次报告：营收+单数+客流+客单价",
        "需求预测：预测营收+置信度+依据",
    ]
    for t in templates:
        pp = tf.add_paragraph()
        pp.text = f"  • {t}"
        pp.font.size = Pt(11)
        pp.font.color.rgb = DARK

    # 右：FSM
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(3.5))
    tf2 = right_box.text_frame
    tf2.word_wrap = True
    p0 = tf2.paragraphs[0]
    p0.text = "行动状态机（FSM）"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    states = "创建 → 推送 → 确认 → 处理中 → 已解决"
    p1 = tf2.add_paragraph()
    p1.text = states
    p1.font.size = Pt(12)
    p1.font.color.rgb = DARK

    # 升级机制
    esc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.0), Inches(9), Inches(1.8))
    esc.fill.solid()
    esc.fill.fore_color.rgb = LIGHT_BG
    esc.line.color.rgb = BRAND
    etf = esc.text_frame
    etf.word_wrap = True
    p0 = etf.paragraphs[0]
    p0.text = "自动升级机制"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    levels = [
        "P0（紧急）：30分钟无响应 → 自动升级到店长",
        "P1（高）  ：2小时无响应 → 升级",
        "P2（中）  ：24小时无响应 → 升级",
        "P3（低）  ：3天无响应 → 升级",
    ]
    for lv in levels:
        pp = etf.add_paragraph()
        pp.text = f"  {lv}"
        pp.font.size = Pt(11)
        pp.font.color.rgb = DARK

    add_footer(slide, 9)


def slide_10_forecast(prs):
    """需求预测三级降级"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "需求预测 — 三级降级策略", 28, True, DARK)

    headers = ["数据天数", "方法", "置信度", "基础"]
    data = [
        ["< 14天", "规则引擎", "低", "行业基准¥3000/天"],
        ["14-60天", "统计模型", "中", "加权移动平均"],
        ["≥ 60天", "Prophet ML", "高", "时间序列预测"],
    ]
    table = make_table(slide, 4, 4, Inches(1), Inches(1.3), Inches(8), Inches(2.0))
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.alignment = PP_ALIGN.CENTER

    # POC策略
    content = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3))
    tf = content.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "POC启动策略"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    steps = [
        "前2周：规则引擎自动启动（无需历史数据）",
        "第3-8周：积累数据后自动切换统计模型",
        "第2月后：数据充足时切换 Prophet ML",
        "",
        "BOM联动：预测营收 → 菜品估算 → 食材用量 × (1+损耗系数)",
    ]
    for s in steps:
        pp = tf.add_paragraph()
        pp.text = f"  {s}" if s else ""
        pp.font.size = Pt(12)
        pp.font.color.rgb = DARK

    add_footer(slide, 10)


def slide_11_hq(prs):
    """多店对标"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "多店对标 — HQ总部看板", 28, True, DARK)

    headers = ["门店", "健康分", "成本率", "营收完成率", "状态"]
    data = [
        ["五一广场店", "76", "35.0%", "85%", "⚠ 成本偏高"],
        ["芙蓉路店", "82", "32.5%", "90%", "良好"],
        ["梅溪湖店", "71", "33.8%", "78%", "⚠ 营收不足"],
        ["株洲店", "68", "34.2%", "72%", "⚠ 多项预警"],
        ["湘潭店", "74", "31.5%", "80%", "良好"],
    ]
    table = make_table(slide, 6, 5, Inches(0.5), Inches(1.3), Inches(9), Inches(3.0))
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                if c == 4 and "⚠" in val:
                    p.font.color.rgb = RED

    # BFF说明
    note = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
    ntf = note.text_frame
    ntf.word_wrap = True
    p0 = ntf.paragraphs[0]
    p0.text = "BFF聚合特性"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = BRAND
    features = [
        "单次请求加载全部门店数据（asyncio.gather并行）",
        "30秒Redis缓存，?refresh=true强制刷新",
        "子调用失败 → 降级返回null，不阻塞整屏",
    ]
    for f in features:
        pp = ntf.add_paragraph()
        pp.text = f"  • {f}"
        pp.font.size = Pt(11)
        pp.font.color.rgb = DARK

    add_footer(slide, 11)


def slide_12_scenarios(prs):
    """POC验证6场景"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "POC验证 — 6大场景", 28, True, DARK)

    headers = ["时段", "场景", "验证Agent", "预期效果"]
    data = [
        ["07:00-09:00", "早晨采购与库存", "Inventory", "采购效率+30%"],
        ["11:30-13:00", "午餐排班优化", "Schedule", "排班效率+40%"],
        ["12:00-12:30", "点餐出餐协同", "Order", "出餐速度+20%"],
        ["12:45", "客诉智能处理", "Service", "处理时间-50%"],
        ["16:00", "晚市销售预测", "Decision", "食材损耗-15%"],
        ["22:30-23:00", "闭店盘点日报", "Performance", "人工成本-10%"],
    ]
    table = make_table(slide, 7, 4, Inches(0.5), Inches(1.3), Inches(9), Inches(4.0))
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                if c == 3:
                    p.font.color.rgb = GREEN
                    p.font.bold = True

    add_footer(slide, 12)


def slide_13_roi(prs):
    """¥价值量化"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "¥价值量化 — 成本节省测算", 28, True, DARK)

    # 单店
    headers = ["节省来源", "月均节省", "年化节省", "置信度"]
    data = [
        ["食材成本率降2点", "¥10,000", "¥120,000", "高"],
        ["损耗减少15%", "¥5,000", "¥60,000", "高"],
        ["人工成本降10%", "¥8,000", "¥96,000", "中"],
        ["采购议价提升", "¥3,000", "¥36,000", "中"],
        ["合计", "¥26,000", "¥312,000", "—"],
    ]
    table = make_table(slide, 6, 4, Inches(0.5), Inches(1.2), Inches(5), Inches(3.0))
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.CENTER
                if r == 5:
                    p.font.bold = True
                    p.font.color.rgb = BRAND

    # ROI大数字
    roi_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(1.2), Inches(3.5), Inches(3.0))
    roi_box.fill.solid()
    roi_box.fill.fore_color.rgb = BRAND
    roi_box.line.fill.background()
    rtf = roi_box.text_frame
    rtf.word_wrap = True
    p0 = rtf.paragraphs[0]
    p0.text = "ROI"
    p0.font.size = Pt(14)
    p0.font.color.rgb = WHITE
    p0.alignment = PP_ALIGN.CENTER
    p1 = rtf.add_paragraph()
    p1.text = "> 500%"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p2 = rtf.add_paragraph()
    p2.text = "5店保守年省"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0xFF, 0xDD, 0xCC)
    p2.alignment = PP_ALIGN.CENTER
    p3 = rtf.add_paragraph()
    p3.text = "¥167万"
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER

    add_footer(slide, 13)


def slide_14_roadmap(prs):
    """实施路线图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    set_text(title, "实施路线图", 28, True, DARK)

    phases = [
        ("Week 1-2", "数据接入 + BOM录入", "规则引擎启动", RGBColor(0x34, 0x98, 0xDB)),
        ("Week 3-4", "健康指数上线 + 企微推送", "第一条AI决策", GREEN),
        ("Week 5-6", "成本分析 + 决策引擎", "验证节省效果", BRAND),
        ("Week 7-8", "统计模型 + 多店对标", "完整POC报告", RGBColor(0x8E, 0x44, 0xAD)),
        ("Month 3+", "ML预测 + 案例积累", "持续优化", DARK),
    ]
    for i, (period, task, milestone, color) in enumerate(phases):
        top = Inches(1.3 + i * 1.15)
        # 时间轴圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), top, Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        # 连线（除最后）
        if i < 4:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.12), Inches(1.3 + i * 1.15 + 0.3), Inches(0.06), Inches(0.85))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            line.line.fill.background()
        # 文字
        tb = slide.shapes.add_textbox(Inches(1.6), Inches(1.2 + i * 1.15), Inches(3), Inches(0.6))
        set_text(tb, period, 16, True, color)
        tb2 = slide.shapes.add_textbox(Inches(4.5), Inches(1.2 + i * 1.15), Inches(3), Inches(0.3))
        set_text(tb2, task, 12, False, DARK)
        tb3 = slide.shapes.add_textbox(Inches(4.5), Inches(1.5 + i * 1.15), Inches(3), Inches(0.3))
        set_text(tb3, f"里程碑: {milestone}", 10, False, GRAY)

    add_footer(slide, 14)


def slide_15_cta(prs):
    """下一步 & CTA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()
    # 品牌条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.8), Inches(8), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()

    # 标题
    title = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(8), Inches(1.5))
    set_text(title, "下一步", 36, True, WHITE, PP_ALIGN.LEFT)
    add_paragraph(title.text_frame, "8周内证明成本率可降 ≥ 1个点", 20, False, BRAND)

    # 请求
    content = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(3.5))
    tf = content.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "我们需要您的支持："
    p0.font.size = Pt(16)
    p0.font.color.rgb = WHITE
    p0.font.bold = True
    asks = [
        "1. 确认五一广场店为首批接入门店",
        "2. 提供天财商龙POS系统接口权限",
        "3. 指定1名店长对接人（企微接收决策推送）",
        "4. 提供5道招牌菜的BOM配方数据",
    ]
    for a in asks:
        pp = tf.add_paragraph()
        pp.text = f"  {a}"
        pp.font.size = Pt(14)
        pp.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        pp.space_after = Pt(6)

    # 承诺
    promise = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.5))
    set_text(promise, "所有过程数据可追溯 — case_story_generator 自动生成案例报告", 11, False, GRAY, PP_ALIGN.LEFT)

    add_footer(slide, 15)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_01_cover(prs)
    slide_02_profile(prs)
    slide_03_value(prs)
    slide_04_architecture(prs)
    slide_05_mvp(prs)
    slide_06_food_cost(prs)
    slide_07_health(prs)
    slide_08_decision(prs)
    slide_09_wechat(prs)
    slide_10_forecast(prs)
    slide_11_hq(prs)
    slide_12_scenarios(prs)
    slide_13_roi(prs)
    slide_14_roadmap(prs)
    slide_15_cta(prs)

    prs.save(OUTPUT)
    print(f"✅ PPT已生成: {OUTPUT}")
    print(f"   共 {len(prs.slides)} 页")


if __name__ == "__main__":
    main()
